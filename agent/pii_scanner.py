import os
import re

# ── Regex patterns ────────────────────────────────────────────────────────────

_IIN_RE = re.compile(r"(?<!\d)(\d{12})(?!\d)")

_PHONE_RE = re.compile(
    r"(?:\+7|8)[\s\-]?\(?\d{3}\)?[\s\-]?\d{3}[\s\-]?\d{2}[\s\-]?\d{2}"
)

_EMAIL_RE = re.compile(r"[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+")

_DOB_RE = re.compile(
    r"\b(?:\d{2}[./]\d{2}[./]\d{4}|\d{4}-\d{2}-\d{2})\b"
)

_FIO_RU_RE = re.compile(
    r"\b[А-ЯЁ][а-яё]+\s+[А-ЯЁ][а-яё]+(?:\s+[А-ЯЁ][а-яё]+)?\b"
)

_FIO_KZ_SUFFIX_RE = re.compile(
    r"\b[А-ЯЁ][а-яёА-ЯЁ]+\s+[А-ЯЁ][а-яёА-ЯЁ]+(?:улы|қызы|ұлы|қызы)\b"
)

# Plain-text extensions scanned directly
_TEXT_EXTS = {".txt", ".csv", ".log", ".json", ".xml", ".tsv", ".md"}

# Office extensions scanned via their respective libraries
_OFFICE_EXTS = {".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".pdf"}

_SUPPORTED_EXTS = _TEXT_EXTS | _OFFICE_EXTS


# ── IIN checksum validation ───────────────────────────────────────────────────

def _validate_iin(iin_str: str) -> bool:
    if len(iin_str) != 12 or not iin_str.isdigit():
        return False
    d  = [int(c) for c in iin_str]
    w1 = [1, 2, 3, 4, 5, 6, 7, 8, 9, 10, 11]
    s1 = sum(d[i] * w1[i] for i in range(11)) % 11
    if s1 < 10:
        return d[11] == s1
    w2 = [3, 4, 5, 6, 7, 8, 9, 10, 11, 1, 2]
    s2 = sum(d[i] * w2[i] for i in range(11)) % 11
    return d[11] == s2


# ── spaCy (optional) ──────────────────────────────────────────────────────────

try:
    import spacy
    _nlp = spacy.load("ru_core_news_sm")
    _SPACY_AVAILABLE = True
except Exception:
    _nlp = None
    _SPACY_AVAILABLE = False


def _extract_spacy_persons(text: str) -> list:
    if not _SPACY_AVAILABLE or not _nlp:
        return []
    doc = _nlp(text[:10_000])
    return [ent.text for ent in doc.ents if ent.label_ == "PER"]


# ── Office text extraction ────────────────────────────────────────────────────

def _cell_str(value) -> str:
    """Convert a spreadsheet cell to string, preserving integer IINs."""
    if isinstance(value, float) and value == int(value):
        return str(int(value))
    return str(value) if value is not None else ""


def _extract_text(path: str, ext: str) -> str:
    """Return plain text from a file regardless of format."""
    try:
        if ext in (".docx", ".doc"):
            from docx import Document
            doc   = Document(path)
            parts = []
            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        for p in cell.paragraphs:
                            if p.text.strip():
                                parts.append(p.text)
            for section in doc.sections:
                for hdr in (section.header, section.footer):
                    if hdr:
                        for p in hdr.paragraphs:
                            if p.text.strip():
                                parts.append(p.text)
            return "\n".join(parts)

        if ext in (".xlsx", ".xls"):
            import openpyxl
            wb   = openpyxl.load_workbook(path, read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    line = " ".join(_cell_str(c) for c in row if c is not None)
                    if line.strip():
                        rows.append(line)
            wb.close()
            return "\n".join(rows)

        if ext in (".pptx", ".ppt"):
            from pptx import Presentation
            prs   = Presentation(path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                if cell.text.strip():
                                    parts.append(cell.text)
            return "\n".join(parts)

        if ext == ".pdf":
            import pdfplumber
            parts = []
            with pdfplumber.open(path) as pdf:
                for page in pdf.pages[:30]:   # limit to first 30 pages
                    t = page.extract_text()
                    if t:
                        parts.append(t)
                        if len("\n".join(parts)) > 500_000:
                            break
            return "\n".join(parts)

    except Exception:
        pass

    # Plain text fallback
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read(500_000)
    except Exception:
        return ""


# ── Main scanner class ────────────────────────────────────────────────────────

class PIIScanner:
    def scan_text(self, text: str) -> tuple:
        """Return (is_pii, confidence, pii_types)."""
        found_types: list = []

        for match in _IIN_RE.finditer(text):
            if _validate_iin(match.group(1)):
                found_types.append("IIN")
                break

        if _PHONE_RE.search(text):
            found_types.append("PHONE")

        if _EMAIL_RE.search(text):
            found_types.append("EMAIL")

        if _DOB_RE.search(text):
            found_types.append("DOB")

        persons = _extract_spacy_persons(text)
        if persons:
            found_types.append("FIO")
        elif _FIO_RU_RE.search(text) or _FIO_KZ_SUFFIX_RE.search(text):
            found_types.append("FIO")

        if not found_types:
            return False, 0.0, []

        weights    = {"IIN": 0.40, "PHONE": 0.20, "EMAIL": 0.15, "FIO": 0.15, "DOB": 0.10}
        confidence = round(min(sum(weights.get(t, 0) for t in found_types), 1.0), 2)
        return True, confidence, found_types

    def scan_file(self, path: str, timeout: float = 30.0) -> tuple:
        """Return (is_pii, confidence, pii_types). Supports text and Office formats."""
        import concurrent.futures

        if not os.path.isfile(path):
            return False, 0.0, []

        ext = os.path.splitext(path)[1].lower()
        if ext not in _SUPPORTED_EXTS:
            return False, 0.0, []

        # Skip very large files (> 50 MB) to avoid multi-minute scans
        try:
            if os.path.getsize(path) > 50 * 1024 * 1024:
                return False, 0.0, []
        except Exception:
            return False, 0.0, []

        def _do_scan():
            if ext in _TEXT_EXTS:
                content = ""
                for enc in ("utf-8-sig", "utf-8", "cp1251", "latin-1"):
                    try:
                        with open(path, "r", encoding=enc) as f:
                            content = f.read(500_000)
                        break
                    except (UnicodeDecodeError, LookupError):
                        continue
                return self.scan_text(content)
            # Office / PDF
            try:
                content = _extract_text(path, ext)
                return self.scan_text(content)
            except Exception:
                return False, 0.0, []

        ex = concurrent.futures.ThreadPoolExecutor(max_workers=1)
        fut = ex.submit(_do_scan)
        try:
            return fut.result(timeout=timeout)
        except concurrent.futures.TimeoutError:
            return False, 0.0, []
        except Exception:
            return False, 0.0, []
        finally:
            ex.shutdown(wait=False)  # don't block on hung PDF threads
